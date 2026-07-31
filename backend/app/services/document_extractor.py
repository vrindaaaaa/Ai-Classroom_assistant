"""
document_extractor.py
=====================
Responsible for extracting raw text from uploaded study documents.

Supported formats
-----------------
- PDF  : via PyMuPDF (fitz)
- DOCX : via python-docx
- PPTX : via python-pptx

Each format has a dedicated function so callers can import them individually.
The public entry-point is ``extract_text(file_path)`` which auto-detects the
format from the file extension and delegates to the correct parser.

Error handling
--------------
All extraction failures are wrapped in ``DocumentExtractionError`` so the
caller only has to catch a single exception type.
"""
from __future__ import annotations

import logging
import re
from pathlib import Path
from typing import Union

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Optional dependencies – each is guarded so unit tests can mock them out
# without actually installing the libraries.
# ---------------------------------------------------------------------------
try:
    import fitz  # PyMuPDF
except ImportError:  # pragma: no cover
    fitz = None

try:
    from docx import Document as DocxDocument
except ImportError:  # pragma: no cover
    DocxDocument = None

try:
    from pptx import Presentation
except ImportError:  # pragma: no cover
    Presentation = None


# ---------------------------------------------------------------------------
# Public exception
# ---------------------------------------------------------------------------
class DocumentExtractionError(Exception):
    """Raised when a document cannot be parsed or contains no extractable text."""


# ---------------------------------------------------------------------------
# Supported extensions registry
# ---------------------------------------------------------------------------
SUPPORTED_EXTENSIONS = {"pdf", "docx", "pptx"}


# ---------------------------------------------------------------------------
# Internal helper – text cleaning
# ---------------------------------------------------------------------------
def _clean_text(raw: str) -> str:
    """Normalize whitespace and remove control characters from extracted text."""
    if not raw:
        return ""
    # Replace various whitespace sequences with a single space, preserve newlines
    text = re.sub(r"[ \t]+", " ", raw)
    # Collapse more than two consecutive newlines into two
    text = re.sub(r"\n{3,}", "\n\n", text)
    # Strip leading/trailing whitespace from every line
    lines = [line.strip() for line in text.splitlines()]
    return "\n".join(lines).strip()


# ---------------------------------------------------------------------------
# PDF extraction
# ---------------------------------------------------------------------------
def extract_pdf_text(file_path: Union[str, Path]) -> str:
    """Extract and clean text from a PDF file using PyMuPDF.

    Parameters
    ----------
    file_path:
        Absolute or relative path to a ``.pdf`` file.

    Returns
    -------
    str
        Cleaned text extracted from all pages.

    Raises
    ------
    DocumentExtractionError
        If PyMuPDF is not installed, the file is unreadable, or extraction fails.
    """
    if fitz is None:
        raise DocumentExtractionError(
            "PyMuPDF is not installed. Run: pip install pymupdf"
        )

    logger.info("[extractor] opening PDF file=%s", file_path)
    try:
        doc = fitz.open(str(file_path))
    except Exception as exc:
        logger.error("[extractor] failed to open PDF file=%s error=%s", file_path, exc)
        raise DocumentExtractionError(
            f"Unable to open PDF file '{file_path}': {exc}"
        ) from exc

    try:
        pages = []
        for page in doc:
            page_text = page.get_text()
            if page_text and page_text.strip():
                pages.append(page_text)
        raw = "\n".join(pages)
        cleaned = _clean_text(raw)
        logger.info("[extractor] PDF extraction succeeded file=%s pages=%d chars=%d", file_path, len(pages), len(cleaned))
        return cleaned
    except Exception as exc:
        logger.error("[extractor] failed to extract PDF file=%s error=%s", file_path, exc)
        raise DocumentExtractionError(
            f"Failed to extract text from PDF '{file_path}': {exc}"
        ) from exc
    finally:
        doc.close()


# ---------------------------------------------------------------------------
# DOCX extraction
# ---------------------------------------------------------------------------
def extract_docx_text(file_path: Union[str, Path]) -> str:
    """Extract and clean text from a DOCX file using python-docx.

    Parameters
    ----------
    file_path:
        Absolute or relative path to a ``.docx`` file.

    Returns
    -------
    str
        Cleaned text extracted from all paragraphs.

    Raises
    ------
    DocumentExtractionError
        If python-docx is not installed, the file is unreadable, or extraction fails.
    """
    if DocxDocument is None:
        raise DocumentExtractionError(
            "python-docx is not installed. Run: pip install python-docx"
        )

    logger.info("[extractor] opening DOCX file=%s", file_path)
    try:
        doc = DocxDocument(str(file_path))
    except Exception as exc:
        logger.error("[extractor] failed to open DOCX file=%s error=%s", file_path, exc)
        raise DocumentExtractionError(
            f"Unable to open DOCX file '{file_path}': {exc}"
        ) from exc

    paragraphs = [
        para.text.strip()
        for para in doc.paragraphs
        if para.text and para.text.strip()
    ]
    raw = "\n".join(paragraphs)
    cleaned = _clean_text(raw)
    logger.info("[extractor] DOCX extraction succeeded file=%s paragraphs=%d chars=%d", file_path, len(paragraphs), len(cleaned))
    return cleaned


# ---------------------------------------------------------------------------
# PPTX extraction
# ---------------------------------------------------------------------------
def extract_pptx_text(file_path: Union[str, Path]) -> str:
    """Extract and clean text from a PPTX file using python-pptx.

    Iterates over every slide and every text-bearing shape to build a
    slide-by-slide transcript.

    Parameters
    ----------
    file_path:
        Absolute or relative path to a ``.pptx`` file.

    Returns
    -------
    str
        Cleaned text extracted from all slides.

    Raises
    ------
    DocumentExtractionError
        If python-pptx is not installed, the file is unreadable, or extraction fails.
    """
    if Presentation is None:
        raise DocumentExtractionError(
            "python-pptx is not installed. Run: pip install python-pptx"
        )

    logger.info("[extractor] opening PPTX file=%s", file_path)
    try:
        prs = Presentation(str(file_path))
    except Exception as exc:
        logger.error("[extractor] failed to open PPTX file=%s error=%s", file_path, exc)
        raise DocumentExtractionError(
            f"Unable to open PPTX file '{file_path}': {exc}"
        ) from exc

    slide_blocks: list[str] = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_lines: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text and shape.text.strip():
                slide_lines.append(shape.text.strip())
        if slide_lines:
            slide_blocks.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_lines))

    raw = "\n\n".join(slide_blocks)
    cleaned = _clean_text(raw)
    logger.info("[extractor] PPTX extraction succeeded file=%s slides=%d chars=%d", file_path, len(slide_blocks), len(cleaned))
    return cleaned


# ---------------------------------------------------------------------------
# Unified entry-point
# ---------------------------------------------------------------------------
def extract_text(file_path: Union[str, Path]) -> str:
    """Detect the file extension and extract text with the appropriate parser.

    Parameters
    ----------
    file_path:
        Path to the uploaded document (PDF, DOCX, or PPTX).

    Returns
    -------
    str
        Clean extracted text ready for AI processing.

    Raises
    ------
    DocumentExtractionError
        If the extension is unsupported or the underlying parser fails.
    """
    path = Path(file_path)
    extension = path.suffix.lower().lstrip(".")

    if not extension:
        raise DocumentExtractionError("The file has no extension; cannot determine type.")

    if extension not in SUPPORTED_EXTENSIONS:
        raise DocumentExtractionError(
            f"Unsupported file type '.{extension}'. "
            f"Supported types: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
        )

    logger.info("[extractor] extracting text from file=%s extension=%s", file_path, extension)
    if extension == "pdf":
        return extract_pdf_text(path)
    if extension == "docx":
        return extract_docx_text(path)
    if extension == "pptx":
        return extract_pptx_text(path)

    # Should never reach here, but keeps mypy happy
    raise DocumentExtractionError(f"Unsupported file type: .{extension}")
